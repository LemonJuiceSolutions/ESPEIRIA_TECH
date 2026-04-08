from __future__ import annotations

import argparse
import os
import subprocess
import sys
from html.parser import HTMLParser
from pathlib import Path


def count_slides(html: str) -> int:
    class SlideCounter(HTMLParser):
        def __init__(self) -> None:
            super().__init__()
            self.count = 0

        def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
            if tag not in {"div", "section"}:
                return
            class_attr = None
            for k, v in attrs:
                if k == "class":
                    class_attr = v or ""
                    break
            if not class_attr:
                return
            tokens = class_attr.split()
            if "slide" in tokens:
                self.count += 1

    p = SlideCounter()
    p.feed(html)
    return p.count


def chrome_path() -> str:
    candidates = [
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "/Applications/Chromium.app/Contents/MacOS/Chromium",
    ]
    for c in candidates:
        if Path(c).exists():
            return c
    raise FileNotFoundError("Chrome/Chromium not found in /Applications")


def ensure_deps(root: Path) -> None:
    try:
        import PIL
        import img2pdf
        import pptx
        return
    except ModuleNotFoundError:
        pass

    if os.environ.get("ESPEIRIA_EXPORT_VENV") == "1":
        raise SystemExit(
            "Dipendenze mancanti. Prova a ricreare il venv:\n"
            "rm -rf .venv_export\n"
            "python3 export.py"
        )

    venv_dir = root / ".venv_export"
    py = venv_dir / "bin" / "python3"
    if not py.exists():
        subprocess.run([sys.executable, "-m", "venv", venv_dir.as_posix()], check=True)

    subprocess.run([py.as_posix(), "-m", "pip", "install", "--upgrade", "pip"], check=True)
    subprocess.run(
        [py.as_posix(), "-m", "pip", "install", "img2pdf", "pillow", "python-pptx"],
        check=True,
    )

    env = os.environ.copy()
    env["ESPEIRIA_EXPORT_VENV"] = "1"
    os.execve(
        py.as_posix(),
        [py.as_posix(), str(Path(__file__).resolve()), *sys.argv[1:]],
        env,
    )


def run() -> int:
    root = Path(__file__).resolve().parent
    ensure_deps(root)
    import img2pdf
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    parser = argparse.ArgumentParser()
    parser.add_argument("--width", type=int, default=1920)
    parser.add_argument("--height", type=int, default=1080)
    parser.add_argument("--capture-height", type=int, default=1165)
    parser.add_argument("--time-budget", type=int, default=9000)
    parser.add_argument("--timeout", type=int, default=45)
    args = parser.parse_args()

    base_name = "ESPEIIRIA_TECH"
    html_path = root / "index.html"
    html = html_path.read_text(encoding="utf-8", errors="ignore")
    n = count_slides(html)
    if n <= 0:
        raise RuntimeError("No slides found")

    out_dir = root / ".export"
    out_dir.mkdir(exist_ok=True)

    chrome = chrome_path()
    base_url = f"file://{html_path.as_posix()}"

    raw_paths: list[Path] = []
    slide_paths: list[Path] = []
    for i in range(1, n + 1):
        print(f"Render {i}/{n}")
        raw = out_dir / f"raw-{i:02d}.png"
        url = f"{base_url}?export=1&slide={i}"
        subprocess.run(
            [
                chrome,
                "--headless=new",
                "--disable-gpu",
                "--hide-scrollbars",
                f"--window-size={args.width},{args.capture_height}",
                f"--virtual-time-budget={args.time_budget}",
                f"--screenshot={raw.as_posix()}",
                url,
            ],
            check=True,
            timeout=args.timeout,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        raw_paths.append(raw)

        with Image.open(raw) as im:
            img = im.convert("RGB").crop((0, 0, args.width, args.height))
            slide = out_dir / f"slide-{i:02d}.png"
            img.save(slide)
        slide_paths.append(slide)

    pdf_path = root / f"{base_name}.pdf"
    with open(pdf_path, "wb") as f:
        f.write(img2pdf.convert([p.as_posix() for p in slide_paths]))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for p in slide_paths:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(p.as_posix(), 0, 0, width=prs.slide_width, height=prs.slide_height)
    pptx_path = root / f"{base_name}.pptx"
    prs.save(pptx_path.as_posix())

    return 0


if __name__ == "__main__":
    raise SystemExit(run())
